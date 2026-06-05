"""Quick inspector for the FFF master template.

Dumps slide masters, layouts, placeholders, and existing slides so we can
target the right layouts when generating the monthly deck.

Run: python3 inspect_template.py
"""

from pptx import Presentation
from pptx.util import Emu

TEMPLATE = "/Users/phil/GitHub/sigmaquickstarts/.claude/SKILLS/first-friday-features/2026 (Template).pptx"


def emu_to_in(emu: int) -> float:
    return round(emu / 914400, 2)


def main() -> None:
    prs = Presentation(TEMPLATE)

    print(f"=== Slide dimensions: {emu_to_in(prs.slide_width)}\" x {emu_to_in(prs.slide_height)}\" ===\n")

    print(f"=== Slide masters: {len(prs.slide_masters)} ===")
    for mi, master in enumerate(prs.slide_masters):
        print(f"\nMaster {mi}: layouts = {len(master.slide_layouts)}")
        for li, layout in enumerate(master.slide_layouts):
            print(f"  Layout {li}: {layout.name!r}")
            for ph in layout.placeholders:
                print(
                    f"    placeholder idx={ph.placeholder_format.idx} "
                    f"type={ph.placeholder_format.type} name={ph.name!r}"
                )

    print(f"\n=== Existing slides in template: {len(prs.slides)} ===")
    for si, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name
        print(f"\nSlide {si + 1}: layout={layout_name!r}")
        for shape in slide.shapes:
            shape_kind = shape.shape_type
            is_ph = shape.is_placeholder
            ph_idx = shape.placeholder_format.idx if is_ph else None
            ph_name = shape.name
            text = ""
            if shape.has_text_frame:
                text = " | ".join(
                    p.text for p in shape.text_frame.paragraphs if p.text
                )[:120]
            print(
                f"  shape name={ph_name!r} kind={shape_kind} ph_idx={ph_idx} text={text!r}"
            )


if __name__ == "__main__":
    main()
