"""FFF monthly slide-deck generator.

Reads an FFF monthly markdown file and produces a .pptx deck by cloning
slides from `Template_first_friday_features.pptx` and substituting content.
Output goes to `output/{MM}_{YYYY}_first_friday_features.pptx`.

Strategy: clone reference slides from the template (title, contents,
spotlight, quick-hits, new-QSes, bug-fixes, closing) and substitute text.
This preserves the template's visual treatment (image column on
spotlights, status-label text boxes, picture/table placement on
contents and closing slides).

Spotlight detection is automatic: any feature heading marked with
`<img src="assets/heart_icon.png" ...>` becomes a spotlight slide,
including features inside the Overview section.

Slide order: Title → Contents → for each category in markdown order
{Spotlights, then Quick Hits if any non-spotlight features remain}
→ New QuickStarts → Bug Fixes → Closing.

Usage:
    python3 generate_deck.py <markdown_path>
"""

from __future__ import annotations

import re
import sys
from copy import deepcopy
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation

# --- paths ---------------------------------------------------------------

SKILL_DIR = Path(__file__).parent
TEMPLATE = SKILL_DIR / "Template_First_Friday_Features.pptx"
OUTPUT_DIR = SKILL_DIR / "output"

# Reference slide indices in the template
REF_TITLE = 0
REF_CONTENTS = 1
REF_SPOTLIGHT = 2
REF_QUICK_HITS = 3
REF_NEW_QS = 4
REF_BUG_FIXES = 5
REF_CLOSING = 6


# --- parsing -------------------------------------------------------------

@dataclass
class Feature:
    title: str
    title_clean: str        # title with status suffix and heart stripped
    category: str
    status: str             # "GA" | "Public Beta" | "Beta" | ""
    body: list[str] = field(default_factory=list)  # description paragraphs
    bullets: list[str] = field(default_factory=list)
    why_it_matters: str = ""
    doc_url: str = ""
    is_heart: bool = False
    related_quickstarts: list[dict] = field(default_factory=list)  # [{title, url}]


@dataclass
class NewQuickStart:
    title: str
    url: str
    description: str = ""


@dataclass
class ParsedFFF:
    month: str
    year: str
    features: list[Feature] = field(default_factory=list)
    new_quickstarts: list[NewQuickStart] = field(default_factory=list)
    bug_fixes: list[str] = field(default_factory=list)


MONTHS = {
    "01": "January", "02": "February", "03": "March", "04": "April",
    "05": "May", "06": "June", "07": "July", "08": "August",
    "09": "September", "10": "October", "11": "November", "12": "December",
}

STATUS_RE = re.compile(r"\((GA|Public Beta|Beta)\)", re.IGNORECASE)
HEART_RE = re.compile(r'<img\s+src=["\']assets/heart_icon\.png["\'][^>]*>')
DOC_LINK_RE = re.compile(
    r"For more information, see \[([^\]]+)\]\(([^)]+)\)",
    re.IGNORECASE,
)
NEW_QS_LINK_RE = re.compile(r"\[This QuickStart\]\(([^)]+)\)")
BUG_FIX_RE = re.compile(r"^\*\*(\d+):\*\*\s*(.+)$")
RELATED_QS_RE = re.compile(
    r"\[([^\]]+)\]\((https://quickstarts\.sigmacomputing\.com[^)]+)\)"
)


def _strip_title(raw: str) -> tuple[str, str, bool]:
    """Return (clean_title, status, is_heart)."""
    is_heart = bool(HEART_RE.search(raw))
    raw = HEART_RE.sub("", raw).strip()
    m = STATUS_RE.search(raw)
    status = ""
    if m:
        status = m.group(1)
        if status.lower() == "beta":
            status = "Beta"
        elif status.lower() == "public beta":
            status = "Public Beta"
        else:
            status = "GA"
        raw = STATUS_RE.sub("", raw).strip()
    return raw, status, is_heart


def parse_markdown(path: Path) -> ParsedFFF:
    text = path.read_text(encoding="utf-8")
    lines = text.splitlines()

    # Month/year from filename pattern MM_YYYY_first_friday_features.md
    fname = path.stem  # e.g., 04_2026_first_friday_features
    m = re.match(r"(\d{2})_(\d{4})_first_friday_features", fname)
    if not m:
        raise SystemExit(f"Filename does not match MM_YYYY pattern: {fname}")
    month = MONTHS[m.group(1)]
    year = m.group(2)

    parsed = ParsedFFF(month=month, year=year)

    current_cat: str | None = None
    current_feat: Feature | None = None

    def flush_feature() -> None:
        nonlocal current_feat
        if current_feat is not None:
            parsed.features.append(current_feat)
        current_feat = None

    i = 0
    while i < len(lines):
        line = lines[i]
        # Category headers
        if line.startswith("## "):
            flush_feature()
            cat_name = line[3:].strip()
            if cat_name.lower().startswith("overview"):
                # Overview can contain heart-tagged Spotlight features (e.g.,
                # platform announcements). Treat as a category so the parser
                # picks them up; non-spotlight ### headings in Overview are
                # filtered out in build_deck.
                current_cat = "Overview"
            elif cat_name.lower().startswith("additional"):
                current_cat = None
            elif cat_name == "Bug Fixes":
                current_cat = "Bug Fixes"
            elif cat_name.startswith("New QuickStarts"):
                current_cat = "New QuickStarts"
            else:
                current_cat = cat_name
            i += 1
            continue

        # Bug fixes — supports two formats:
        #   (a) `**N:** description text` (April-era numbered style)
        #   (b) `### Bug Title \n description paragraph` (newer titled style)
        if current_cat == "Bug Fixes":
            # Format (a): numbered single-line
            m = BUG_FIX_RE.match(line.strip())
            if m:
                parsed.bug_fixes.append(m.group(2).strip())
                i += 1
                continue
            # Format (b): ### Title followed by description
            if line.startswith("### "):
                bf_title = line[4:].strip()
                bf_desc = ""
                j = i + 1
                while j < len(lines):
                    ln = lines[j].strip()
                    if ln.startswith("###") or ln.startswith("##"):
                        break
                    if ln and not ln.startswith("<") and not ln.startswith("!["):
                        bf_desc = ln
                        break
                    j += 1
                if bf_desc:
                    parsed.bug_fixes.append(f"{bf_title} — {bf_desc}")
                else:
                    parsed.bug_fixes.append(bf_title)
                i += 1
                continue
            i += 1
            continue

        # New QuickStarts — each ### is a QS
        if current_cat == "New QuickStarts" and line.startswith("### "):
            qs_title = line[4:].strip()
            qs_title_clean, _, _ = _strip_title(qs_title)
            qs_url = ""
            qs_desc = ""
            j = i + 1
            while j < len(lines) and not lines[j].startswith("### ") and not lines[j].startswith("## "):
                m = NEW_QS_LINK_RE.search(lines[j])
                if m:
                    qs_url = m.group(1)
                    # The description lives on the same line; replace the
                    # markdown link with plain text "This QuickStart" so the
                    # sentence reads naturally on the slide.
                    qs_desc = re.sub(
                        r"\[This QuickStart\]\([^)]+\)",
                        "This QuickStart",
                        lines[j],
                    ).strip()
                    break
                j += 1
            parsed.new_quickstarts.append(
                NewQuickStart(title=qs_title_clean, url=qs_url, description=qs_desc)
            )
            i += 1
            continue

        # Feature headers — ### within a content category
        if line.startswith("### ") and current_cat and current_cat not in {"Bug Fixes", "New QuickStarts"}:
            flush_feature()
            title_raw = line[4:].strip()
            title_clean, status, is_heart = _strip_title(title_raw)
            current_feat = Feature(
                title=title_raw,
                title_clean=title_clean,
                category=current_cat,
                status=status or "GA",
                is_heart=is_heart,
            )
            i += 1
            continue

        # Within a feature: collect body, bullets, doc link, WIM
        if current_feat is not None:
            stripped = line.strip()

            # WHY IT MATTERS block (markdown bold form)
            if "**WHY IT MATTERS:**" in stripped:
                # Capture subsequent non-blank lines until blank or next ###/##
                wim_parts: list[str] = []
                j = i + 1
                while j < len(lines):
                    ln = lines[j].rstrip()
                    if ln.startswith("###") or ln.startswith("##"):
                        break
                    if ln.strip() == "" and wim_parts:
                        break
                    if ln.strip():
                        wim_parts.append(ln.strip())
                    j += 1
                current_feat.why_it_matters = " ".join(wim_parts)
                i = j
                continue

            # Doc link
            m = DOC_LINK_RE.search(stripped)
            if m and not current_feat.doc_url:
                current_feat.doc_url = m.group(2)
                i += 1
                continue

            # Related QuickStart reference — strip from body and capture
            # separately so it can be rendered as a clickable bullet.
            qs_m = RELATED_QS_RE.search(stripped)
            if qs_m:
                current_feat.related_quickstarts.append(
                    {"title": qs_m.group(1), "url": qs_m.group(2)}
                )
                i += 1
                continue

            # Bullets
            if stripped.startswith("- ") or stripped.startswith("* "):
                current_feat.bullets.append(stripped[2:].strip())
                i += 1
                continue

            # Body text
            if stripped and not stripped.startswith("<") and not stripped.startswith("!["):
                current_feat.body.append(stripped)
            i += 1
            continue

        i += 1

    flush_feature()
    return parsed


# --- slide manipulation --------------------------------------------------

def clone_slide(prs, src_slide):
    """Deep-clone a slide and append at end of presentation."""
    blank = prs.slide_layouts.get_by_name(src_slide.slide_layout.name) \
        if hasattr(prs.slide_layouts, "get_by_name") else src_slide.slide_layout
    new_slide = prs.slides.add_slide(src_slide.slide_layout)
    # Remove default placeholders added by layout (so we use only cloned shapes)
    for sp in list(new_slide.shapes):
        sp.element.getparent().remove(sp.element)
    # Copy all shape elements from source
    for shape in src_slide.shapes:
        el = deepcopy(shape.element)
        new_slide.shapes._spTree.append(el)
    return new_slide


def set_text(shape, text: str, hyperlink: str | None = None) -> None:
    """Replace text in a shape's first paragraph, keeping first-run style."""
    from pptx.oxml.ns import qn

    tf = shape.text_frame
    # Clear paragraphs beyond first
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    p = tf.paragraphs[0]
    # Keep first run for style, clear extras
    runs = list(p.runs)
    if not runs:
        run = p.add_run()
    else:
        run = runs[0]
        for extra in runs[1:]:
            extra._r.getparent().remove(extra._r)
    # Strip any existing hyperlink XML from this run. Cloned slides carry
    # hlinkClick elements whose r:id references the source slide's
    # relationships, which weren't copied — pptx will raise KeyError on
    # drop_rel when assigning a new hyperlink.address otherwise.
    rPr = run._r.find(qn("a:rPr"))
    if rPr is not None:
        for hl in rPr.findall(qn("a:hlinkClick")):
            rPr.remove(hl)
    run.text = text
    if hyperlink:
        run.hyperlink.address = hyperlink


def set_bullets(shape, bullets) -> None:
    """Replace bullets in a body shape.

    `bullets` can be either:
      - list[str]                          → all level-0, no hyperlinks
      - list[dict]                         → dicts with keys {text, level?, hyperlink?}
    """
    from pptx.oxml.ns import qn

    tf = shape.text_frame
    if not bullets:
        tf.paragraphs[0].text = ""
        return

    # Normalize input to dict form
    normalized = []
    for b in bullets:
        if isinstance(b, str):
            normalized.append({"text": b, "level": 0, "hyperlink": None})
        else:
            normalized.append(
                {
                    "text": b.get("text", ""),
                    "level": b.get("level", 0),
                    "hyperlink": b.get("hyperlink"),
                }
            )

    # Keep first paragraph as style template
    base_p = tf.paragraphs[0]
    base_runs = list(base_p.runs)
    base_run_xml = deepcopy(base_runs[0]._r) if base_runs else None
    base_pPr = base_p._p.find(qn("a:pPr"))
    base_pPr_xml = deepcopy(base_pPr) if base_pPr is not None else None

    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)

    for idx, bullet in enumerate(normalized):
        text = bullet["text"]
        level = bullet["level"]
        hyperlink = bullet["hyperlink"]

        if idx == 0:
            p = base_p
        else:
            p = tf.add_paragraph()
            if base_pPr_xml is not None:
                p._p.insert(0, deepcopy(base_pPr_xml))

        # Set paragraph level (sub-bullet)
        if level > 0:
            pPr = p._p.find(qn("a:pPr"))
            if pPr is None:
                pPr = p._p.makeelement(qn("a:pPr"), {})
                p._p.insert(0, pPr)
            pPr.set("lvl", str(level))

        # Clear existing runs
        for r in list(p.runs):
            r._r.getparent().remove(r._r)

        # Add run with base style cloned so size/font carry over
        if base_run_xml is not None:
            r = deepcopy(base_run_xml)
            t = r.find(qn("a:t"))
            if t is not None:
                t.text = text
            p._p.append(r)
            if hyperlink:
                rPr = r.find(qn("a:rPr"))
                if rPr is None:
                    rPr = r.makeelement(qn("a:rPr"), {})
                    r.insert(0, rPr)
                hlink = rPr.makeelement(
                    qn("a:hlinkClick"),
                    {qn("r:id"): _add_rel(p.part, hyperlink)},
                )
                rPr.append(hlink)
        else:
            run = p.add_run()
            run.text = text
            if hyperlink:
                run.hyperlink.address = hyperlink


def _add_rel(part, url: str) -> str:
    """Add an external relationship and return its rId."""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    return part.relate_to(url, RT.HYPERLINK, is_external=True)


def find_shape_by_text(slide, needle: str):
    for shape in slide.shapes:
        if shape.has_text_frame and needle in shape.text_frame.text:
            return shape
    return None


def find_placeholder(slide, idx: int):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def find_body_placeholder(slide):
    """Return the body placeholder for a content slide.

    Tries the standard ph_idx=1 first; falls back to any non-title placeholder
    with a text frame. Some Google-Slides-exported templates use placeholders
    without a meaningful ph_idx (reported as 4294967295 by python-pptx).
    """
    body = find_placeholder(slide, 1)
    if body is not None:
        return body
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            continue
        if ph.has_text_frame:
            return ph
    return None


def find_related_qs_box(slide):
    """Return the dedicated 'QuickStart: ...' text box on a Spotlight slide.

    Phil's template includes a visual element specifically for the Related
    QuickStart marker; this helper locates it so the generator can populate
    or remove it as appropriate.
    """
    for shape in slide.shapes:
        if shape.is_placeholder:
            continue
        if shape.has_text_frame and shape.text_frame.text.strip().startswith("QuickStart:"):
            return shape
    return None


# --- deck construction ---------------------------------------------------

def _build_spotlight_slide(prs, ref_spotlight, feat: "Feature") -> tuple[bool, list[str]]:
    """Clone the Spotlight reference slide and populate it with a feature.

    Returns (had_wim, flagged_issues).
    """
    slide = clone_slide(prs, ref_spotlight)
    title_ph = find_placeholder(slide, 0)
    body_ph = find_body_placeholder(slide)
    related_qs_box = find_related_qs_box(slide)

    status_box = None
    for shape in slide.shapes:
        # Skip the related-QS box (it starts with "QuickStart:") so we don't
        # mistake it for a status label.
        if shape is related_qs_box:
            continue
        if shape.has_text_frame and shape.text_frame.text.strip() in {
            "GA", "Public Beta", "Beta", "{RELEASE STATUS}"
        }:
            status_box = shape
            break

    if title_ph:
        set_text(title_ph, feat.title_clean, hyperlink=feat.doc_url or None)

    issues: list[str] = []
    has_wim = bool(feat.why_it_matters)

    if body_ph:
        bullet_items: list = []
        if feat.bullets:
            bullet_items.extend(feat.bullets)
        elif feat.body:
            bullet_items.extend(feat.body)
        if feat.why_it_matters:
            bullet_items.append(f"Why it matters: {feat.why_it_matters}")
        if not bullet_items:
            bullet_items = [feat.title_clean]
        set_bullets(body_ph, bullet_items)

    # Related-QuickStart text box: populate with the first related QS if there
    # is one, otherwise remove the box so it doesn't show placeholder text.
    if related_qs_box is not None:
        if feat.related_quickstarts:
            ref = feat.related_quickstarts[0]
            set_text(
                related_qs_box,
                f"QuickStart: {ref['title']}",
                hyperlink=ref["url"],
            )
        else:
            related_qs_box.element.getparent().remove(related_qs_box.element)

    if status_box:
        set_text(status_box, feat.status or "GA")

    return has_wim, issues


def _build_quick_hits_slide(prs, ref_quick_hits, slide_title: str,
                            features: list["Feature"]) -> list[str]:
    """Clone the Quick Hits reference slide for a category and populate.

    Returns flagged_no_url list for features lacking a doc link.
    """
    slide = clone_slide(prs, ref_quick_hits)
    title_ph = find_placeholder(slide, 0)
    body_ph = find_body_placeholder(slide)

    if title_ph:
        set_text(title_ph, slide_title)

    flagged: list[str] = []
    has_beta_on_slide = False
    bullet_items: list[dict] = []

    for f in features:
        title_text = f.title_clean
        if f.status in {"Beta", "Public Beta"}:
            title_text = f"{title_text} — Public Beta"
            has_beta_on_slide = True
        bullet_items.append(
            {"text": title_text, "level": 0, "hyperlink": f.doc_url or None}
        )
        if not f.doc_url:
            flagged.append(f"{f.category}: {f.title_clean}")

        # Single-sentence description per the SKILL Quick Hits rule.
        # No separate WHY IT MATTERS block in Quick Hits.
        desc = f.body[0] if f.body else ""
        if desc:
            if len(desc) > 220:
                desc = desc[:217] + "..."
            bullet_items.append({"text": desc, "level": 1})

    if body_ph:
        set_bullets(body_ph, bullet_items)

    # Always remove the slide-level GA label; Public Beta label kept only when
    # at least one feature on the slide is Beta.
    for shape in list(slide.shapes):
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if txt == "GA":
            shape.element.getparent().remove(shape.element)
        elif txt in {"Public Beta", "Beta"} and not has_beta_on_slide:
            shape.element.getparent().remove(shape.element)

    return flagged


def build_deck(parsed: ParsedFFF, template_path: Path, output_path: Path) -> dict:
    prs = Presentation(template_path)

    ref_title = prs.slides[REF_TITLE]
    ref_contents = prs.slides[REF_CONTENTS]
    ref_spotlight = prs.slides[REF_SPOTLIGHT]
    ref_quick_hits = prs.slides[REF_QUICK_HITS]
    ref_new_qs = prs.slides[REF_NEW_QS]
    ref_bug_fixes = prs.slides[REF_BUG_FIXES]
    ref_closing = prs.slides[REF_CLOSING]

    flagged_no_url: list[str] = []
    missing_wim: list[str] = []

    # --- Title slide ---
    # The title placeholder text is left as-is from the template; only the
    # subtitle (Month YYYY) is generated. set_text would mangle the soft
    # line break in the title.
    title_slide = clone_slide(prs, ref_title)
    subtitle_ph = find_placeholder(title_slide, 1)
    if subtitle_ph:
        set_text(subtitle_ph, f"{parsed.month} {parsed.year}")

    # --- Contents slide (cloned verbatim — fixed labels) ---
    clone_slide(prs, ref_contents)

    # --- Spotlights + Quick Hits, grouped by category in markdown order ---
    # Track category order from the parsed features so the deck flows in the
    # same order the markdown introduces categories.
    seen_categories: list[str] = []
    for feat in parsed.features:
        if feat.category not in seen_categories:
            seen_categories.append(feat.category)

    spotlight_count = 0
    qh_count = 0
    for category in seen_categories:
        cat_features = [f for f in parsed.features if f.category == category]
        spotlights = [f for f in cat_features if f.is_heart]
        # Overview only contributes Spotlights — its non-heart ### entries
        # (like "Subscribe to What's New in Sigma") are intro matter, not
        # release features, and never get a Quick Hits slide.
        if category == "Overview":
            non_spotlights = []
        else:
            non_spotlights = [f for f in cat_features if not f.is_heart]

        for feat in spotlights:
            had_wim, _ = _build_spotlight_slide(prs, ref_spotlight, feat)
            if not had_wim:
                missing_wim.append(feat.title_clean)
            if not feat.doc_url:
                flagged_no_url.append(f"{category}: {feat.title_clean}")
            spotlight_count += 1

        if non_spotlights:
            qh_title = f"{category} - Quick Hits"
            flagged_no_url.extend(
                _build_quick_hits_slide(prs, ref_quick_hits, qh_title, non_spotlights)
            )
            qh_count += 1

    # --- New QuickStarts slide ---
    if parsed.new_quickstarts:
        qs_slide = clone_slide(prs, ref_new_qs)
        title_ph = find_placeholder(qs_slide, 0)
        body_ph = find_body_placeholder(qs_slide)
        if title_ph:
            set_text(title_ph, f"New QuickStarts in {parsed.month}")
        if body_ph:
            qs_items: list[dict] = []
            for qs in parsed.new_quickstarts:
                qs_items.append(
                    {"text": qs.title, "level": 0, "hyperlink": qs.url or None}
                )
                if qs.description:
                    qs_items.append({"text": qs.description, "level": 1})
            set_bullets(body_ph, qs_items)

    # --- Bug Fixes slide ---
    if parsed.bug_fixes:
        bf_slide = clone_slide(prs, ref_bug_fixes)
        title_ph = find_placeholder(bf_slide, 0)
        body_ph = find_body_placeholder(bf_slide)
        if title_ph:
            set_text(title_ph, "Bug Fixes")
        if body_ph:
            bullets = []
            for bf in parsed.bug_fixes:
                clean = bf.replace("`", "")
                if len(clean) > 140:
                    clean = clean[:137] + "..."
                bullets.append(clean)
            set_bullets(body_ph, bullets)

    # --- Closing slide (cloned verbatim) ---
    clone_slide(prs, ref_closing)

    # Delete the original 7 template reference slides
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst)[:7]:
        rId = sldId.attrib[
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        ]
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)

    return {
        "slide_count": len(prs.slides),
        "spotlight_count": spotlight_count,
        "qh_count": qh_count,
        "no_doc_url": flagged_no_url,
        "missing_wim": missing_wim,
        "output": str(output_path),
    }


def main(argv: list[str]) -> int:
    if len(argv) < 2:
        print(__doc__)
        return 1
    md_path = Path(argv[1])
    if not md_path.exists():
        print(f"Markdown not found: {md_path}")
        return 1

    parsed = parse_markdown(md_path)
    print(f"Parsed: {parsed.month} {parsed.year}")
    print(f"  features: {len(parsed.features)} "
          f"(hearts: {sum(1 for f in parsed.features if f.is_heart)})")
    print(f"  new QSes: {len(parsed.new_quickstarts)}")
    print(f"  bug fixes: {len(parsed.bug_fixes)}")

    output = OUTPUT_DIR / f"{md_path.stem}.pptx"
    stats = build_deck(parsed, TEMPLATE, output)
    print(f"\nGenerated: {output}")
    print(f"  total slides: {stats['slide_count']}")
    print(f"  spotlight slides: {stats['spotlight_count']}")
    print(f"  quick-hits slides: {stats['qh_count']}")
    if stats["no_doc_url"]:
        print(f"  features with no doc URL ({len(stats['no_doc_url'])}):")
        for line in stats["no_doc_url"]:
            print(f"    - {line}")
    if stats["missing_wim"]:
        print(f"  Spotlight features missing WHY IT MATTERS ({len(stats['missing_wim'])}):")
        for line in stats["missing_wim"]:
            print(f"    - {line}")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
